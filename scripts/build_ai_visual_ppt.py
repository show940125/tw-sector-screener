from __future__ import annotations

import argparse
import csv
from datetime import date, datetime
from pathlib import Path
from statistics import median
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT_DIR = Path(__file__).resolve().parents[1]

PALETTE = {
    "bg_dark": RGBColor(20, 31, 48),
    "bg_light": RGBColor(244, 247, 252),
    "ink_dark": RGBColor(26, 35, 51),
    "ink_muted": RGBColor(85, 97, 116),
    "accent": RGBColor(0, 143, 122),
    "accent_alt": RGBColor(0, 107, 179),
    "warn": RGBColor(201, 76, 76),
    "good": RGBColor(43, 138, 62),
    "card": RGBColor(230, 238, 247),
}

WEIGHTS = {
    "Trend": 35,
    "Momentum": 25,
    "Value": 20,
    "Fundamental": 15,
    "Risk": 5,
}

HORIZONS = [
    ("ret_5d", "5D"),
    ("ret_20d", "20D"),
    ("momentum63", "63D"),
    ("momentum126", "126D"),
]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="AI 類股可視化分析 + PPT 生成")
    parser.add_argument(
        "--csv-path",
        default=r"C:\Users\a0953041880\tw-reports\hot5\sector-top100-20260220\theme-ai-4e44e520d4.csv",
        help="AI 類股 CSV 檔",
    )
    parser.add_argument("--as-of", default="2026-02-20", help="分析日期 YYYY-MM-DD")
    parser.add_argument(
        "--output-dir",
        default=r"C:\Users\a0953041880\tw-reports\hot5",
        help="輸出目錄",
    )
    return parser.parse_args()


def _to_num(series: pd.Series) -> pd.Series:
    return pd.to_numeric(series, errors="coerce")


def _safe_text(value: Any, ndigits: int = 2, suffix: str = "") -> str:
    if value is None:
        return "N/A"
    if isinstance(value, float):
        if pd.isna(value):
            return "N/A"
        return f"{value:.{ndigits}f}{suffix}"
    return f"{value}{suffix}"


def _apply_bg(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_title(slide, title: str, subtitle: str, dark: bool = False) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(34)
    p.font.color.rgb = RGBColor(255, 255, 255) if dark else PALETTE["ink_dark"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12.0), Inches(0.6))
    stf = sub_box.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(16)
    sp.font.color.rgb = RGBColor(220, 230, 245) if dark else PALETTE["ink_muted"]
    sp.font.name = "Calibri"
    sp.alignment = PP_ALIGN.LEFT


def _add_note(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def _add_card(slide, x: float, y: float, w: float, h: float, title: str, value: str, tone: str = "normal") -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill = shape.fill
    fill.solid()
    if tone == "good":
        fill.fore_color.rgb = RGBColor(227, 244, 233)
    elif tone == "warn":
        fill.fore_color.rgb = RGBColor(253, 234, 234)
    else:
        fill.fore_color.rgb = PALETTE["card"]
    shape.line.color.rgb = RGBColor(210, 220, 235)

    tf = shape.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = "Calibri"
    p0.font.bold = True
    p0.font.size = Pt(13)
    p0.font.color.rgb = PALETTE["ink_muted"]
    p1 = tf.add_paragraph()
    p1.text = value
    p1.font.name = "Calibri"
    p1.font.bold = True
    p1.font.size = Pt(24)
    p1.font.color.rgb = PALETTE["ink_dark"]


def _format_table_font(table, size: int = 11, header: bool = True) -> None:
    rows = len(table.rows)
    cols = len(table.columns)
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            tf = cell.text_frame
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(size)
            if header and row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PALETTE["accent_alt"]
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)


def _build_summary(df: pd.DataFrame) -> dict[str, Any]:
    summary: dict[str, Any] = {}
    summary["count"] = len(df)
    summary["score_mean"] = float(df["total_score"].mean())
    summary["score_median"] = float(median(df["total_score"].tolist()))
    summary["score_q75"] = float(df["total_score"].quantile(0.75))
    summary["score_q90"] = float(df["total_score"].quantile(0.90))
    summary["top1"] = f"{df.iloc[0]['symbol']} {df.iloc[0]['name']}"

    market_counts = df["market"].value_counts(dropna=False)
    summary["market_counts"] = market_counts.to_dict()

    industry_counts = df["industry"].value_counts(dropna=False)
    summary["industry_counts"] = industry_counts.to_dict()

    avg_returns = {}
    pos_ratios = {}
    for col, label in HORIZONS:
        avg_returns[label] = float(df[col].mean())
        pos_ratios[label] = float((df[col] > 0).mean() * 100.0)
    summary["avg_returns"] = avg_returns
    summary["pos_ratios"] = pos_ratios

    factor_cols = ["trend_score", "momentum_score", "value_score", "fundamental_score", "risk_control_score"]
    summary["factor_means"] = {c: float(df[c].mean()) for c in factor_cols}
    summary["factor_stds"] = {c: float(df[c].std(ddof=0)) for c in factor_cols}
    summary["vol_median"] = float(df["volatility20"].median())
    summary["rsi_median"] = float(df["rsi14"].median())
    summary["ma_stack_counts"] = df["ma_stack"].value_counts().to_dict()
    return summary


def _write_markdown(path: Path, as_of: date, df: pd.DataFrame, top11: pd.DataFrame, summary: dict[str, Any]) -> None:
    lines = [
        "# AI 類股可視化分析摘要",
        "",
        f"- 分析日期：`{as_of.isoformat()}`",
        f"- 樣本數：`{summary['count']}`",
        f"- 第一名：`{summary['top1']}`",
        f"- 分數均值 / 中位數：`{summary['score_mean']:.2f}` / `{summary['score_median']:.2f}`",
        "",
        "## 產業與市場結構",
        f"- 市場分布：{summary['market_counts']}",
        f"- 產業前五：{dict(list(summary['industry_counts'].items())[:5])}",
        "",
        "## 短中長趨勢（全樣本平均）",
    ]
    for label in ["5D", "20D", "63D", "126D"]:
        lines.append(
            f"- {label}：平均報酬 `{summary['avg_returns'][label]:.2f}%`，上漲比例 `{summary['pos_ratios'][label]:.1f}%`"
        )

    lines.extend(
        [
            "",
            "## 因子分數（平均）",
            f"- Trend `{summary['factor_means']['trend_score']:.2f}`",
            f"- Momentum `{summary['factor_means']['momentum_score']:.2f}`",
            f"- Value `{summary['factor_means']['value_score']:.2f}`",
            f"- Fundamental `{summary['factor_means']['fundamental_score']:.2f}`",
            f"- Risk `{summary['factor_means']['risk_control_score']:.2f}`",
            "",
            "## Top11 清單",
            "| Rank | Symbol | Name | Total | 5D% | 20D% | 63D% | 126D% | RSI14 | Vol20 | MA 結構 |",
            "|---:|---|---|---:|---:|---:|---:|---:|---:|---:|---|",
        ]
    )
    for _, row in top11.iterrows():
        lines.append(
            "| {rank} | {symbol} | {name} | {total:.2f} | {r5:.2f} | {r20:.2f} | {r63:.2f} | {r126:.2f} | {rsi:.2f} | {vol:.2f} | {stack} |".format(
                rank=int(row["rank"]),
                symbol=row["symbol"],
                name=row["name"],
                total=float(row["total_score"]),
                r5=float(row["ret_5d"]),
                r20=float(row["ret_20d"]),
                r63=float(row["momentum63"]),
                r126=float(row["momentum126"]),
                rsi=float(row["rsi14"]),
                vol=float(row["volatility20"]),
                stack=row["ma_stack"],
            )
        )
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _build_ppt(path: Path, as_of: date, df: pd.DataFrame, top11: pd.DataFrame, summary: dict[str, Any]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(slide, PALETTE["bg_dark"])
    _add_title(slide, "AI 類股可視化分析簡報", f"資料日期 {as_of.isoformat()} | Top100 視角", dark=True)
    hero = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(11.8), Inches(2.6))
    htf = hero.text_frame
    htf.clear()
    p = htf.paragraphs[0]
    p.text = "重點：產業結構、短中長趨勢、前 11 名個股趨勢拆解"
    p.font.name = "Calibri"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 240, 255)
    _add_note(slide, "先定義這份 deck 的任務：不是預測，而是把 AI 類股 Top100 的結構與趨勢拆成可行動訊號。")

    # 2. Executive summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(slide, PALETTE["bg_light"])
    _add_title(slide, "一頁總結", "先看盤面結構與勝率，再看個股細節")
    _add_card(slide, 0.7, 2.0, 2.9, 1.5, "樣本數", str(summary["count"]))
    _add_card(slide, 3.9, 2.0, 2.9, 1.5, "總分均值 / 中位數", f"{summary['score_mean']:.2f} / {summary['score_median']:.2f}")
    _add_card(slide, 7.1, 2.0, 2.9, 1.5, "Top1", summary["top1"], tone="good")
    _add_card(slide, 10.3, 2.0, 2.3, 1.5, "RSI 中位數", f"{summary['rsi_median']:.2f}")
    _add_card(slide, 0.7, 3.8, 3.8, 1.5, "20D 平均報酬", f"{summary['avg_returns']['20D']:.2f}%")
    _add_card(slide, 4.8, 3.8, 3.8, 1.5, "63D 平均報酬", f"{summary['avg_returns']['63D']:.2f}%", tone="good")
    _add_card(slide, 8.9, 3.8, 3.7, 1.5, "126D 平均報酬", f"{summary['avg_returns']['126D']:.2f}%", tone="good")
    _add_note(slide, "這一頁先回答兩件事：整體趨勢是否向上、風險是否過熱。")

    # 3. Scoring method + weights
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(slide, RGBColor(255, 255, 255))
    _add_title(slide, "評分方法與權重", "分數是排序器，不是報酬保證")
    body = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(5.4), Inches(3.6))
    btf = body.text_frame
    btf.clear()
    items = [
        "Total = 0.35*Trend + 0.25*Momentum + 0.20*Value + 0.15*Fundamental + 0.05*Risk",
        "Trend：均線結構 + RSI 規則分",
        "Momentum：63/126 日報酬分位",
        "Value：PE/PB 反向分位 + 殖利率分位",
        "Fundamental：營收 YoY/MoM 分位",
        "Risk：低波動 + 高流動性",
    ]
    for idx, text in enumerate(items):
        p = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
        p.text = text
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = PALETTE["ink_dark"]

    chart_data = CategoryChartData()
    chart_data.categories = list(WEIGHTS.keys())
    chart_data.add_series("權重(%)", list(WEIGHTS.values()))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(6.4),
        Inches(1.9),
        Inches(6.2),
        Inches(3.9),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 40
    _add_note(slide, "權重代表決策偏好：目前 60% 放在趨勢+動能，因此更偏順勢策略。")

    # 4. Market + industry structure
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(slide, PALETTE["bg_light"])
    _add_title(slide, "產業與市場結構", "AI 主題母體內的市場分布與子產業集中度")

    market_data = CategoryChartData()
    market_data.categories = list(summary["market_counts"].keys())
    market_data.add_series("檔數", list(summary["market_counts"].values()))
    chart1 = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(0.8),
        Inches(1.8),
        Inches(5.8),
        Inches(4.7),
        market_data,
    ).chart
    chart1.has_legend = True
    chart1.legend.position = XL_LEGEND_POSITION.RIGHT

    top_industries = list(summary["industry_counts"].items())[:8]
    ind_data = CategoryChartData()
    ind_data.categories = [x[0] for x in top_industries]
    ind_data.add_series("檔數", [x[1] for x in top_industries])
    chart2 = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(6.8),
        Inches(1.8),
        Inches(5.8),
        Inches(4.7),
        ind_data,
    ).chart
    chart2.has_legend = False
    _add_note(slide, "看兩件事：主題是否過度集中在單一子產業，以及市場板塊是否單邊。")

    # 5. Short/Mid/Long trends
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(slide, RGBColor(255, 255, 255))
    _add_title(slide, "短中長趨勢", "全樣本平均報酬與上漲比例")

    avg_data = CategoryChartData()
    labels = [label for _, label in HORIZONS]
    avg_data.categories = labels
    avg_data.add_series("平均報酬%", [summary["avg_returns"][label] for label in labels])
    chart3 = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.8),
        Inches(1.9),
        Inches(6.0),
        Inches(4.8),
        avg_data,
    ).chart
    chart3.has_legend = False

    pos_data = CategoryChartData()
    pos_data.categories = labels
    pos_data.add_series("上漲比例%", [summary["pos_ratios"][label] for label in labels])
    chart4 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7.0),
        Inches(1.9),
        Inches(5.6),
        Inches(4.8),
        pos_data,
    ).chart
    chart4.has_legend = False
    _add_note(slide, "如果 5D 回檔但 63D/126D 仍強，通常是趨勢中繼而非結構轉空。")

    # 6. Factor profile
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(slide, PALETTE["bg_light"])
    _add_title(slide, "因子剖面", "AI 母體平均分數與離散程度")

    factor_map = {
        "trend_score": "Trend",
        "momentum_score": "Momentum",
        "value_score": "Value",
        "fundamental_score": "Fundamental",
        "risk_control_score": "Risk",
    }
    factor_data = CategoryChartData()
    factor_data.categories = list(factor_map.values())
    factor_data.add_series("平均分", [summary["factor_means"][k] for k in factor_map.keys()])
    chart5 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(1.8),
        Inches(7.1),
        Inches(4.9),
        factor_data,
    ).chart
    chart5.has_legend = False

    std_table = slide.shapes.add_table(6, 2, Inches(8.2), Inches(1.8), Inches(4.5), Inches(4.9)).table
    std_table.cell(0, 0).text = "因子"
    std_table.cell(0, 1).text = "標準差"
    for idx, (key, name) in enumerate(factor_map.items(), start=1):
        std_table.cell(idx, 0).text = name
        std_table.cell(idx, 1).text = f"{summary['factor_stds'][key]:.2f}"
    _format_table_font(std_table, size=12, header=True)
    _add_note(slide, "平均分看方向，標準差看分化。分化高的因子適合做相對強弱挑選。")

    # 7. Top11 ranking
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(slide, RGBColor(255, 255, 255))
    _add_title(slide, "前 11 名個股總分排名", "以 total_score 由高到低")
    labels = [f"{int(r.symbol)} {r.name}" for r in top11.itertuples(index=False)]
    values = [float(x) for x in top11["total_score"].tolist()]
    rank_data = CategoryChartData()
    rank_data.categories = labels
    rank_data.add_series("總分", values)
    chart6 = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.8),
        Inches(1.7),
        Inches(12.0),
        Inches(5.3),
        rank_data,
    ).chart
    chart6.has_legend = False
    _add_note(slide, "先看梯隊：第一梯與第二梯差距，決定倉位是否集中。")

    # 8. Top11 trend matrix
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(slide, PALETTE["bg_light"])
    _add_title(slide, "前 11 名趨勢矩陣", "短中長報酬 + RSI + 均線結構")
    columns = ["股票", "5D%", "20D%", "63D%", "126D%", "RSI14", "MA"]
    table = slide.shapes.add_table(12, len(columns), Inches(0.5), Inches(1.6), Inches(12.4), Inches(5.6)).table
    for col_idx, name in enumerate(columns):
        table.cell(0, col_idx).text = name
    for row_idx, row in enumerate(top11.itertuples(index=False), start=1):
        table.cell(row_idx, 0).text = f"{int(row.symbol)} {row.name}"
        table.cell(row_idx, 1).text = f"{row.ret_5d:.2f}"
        table.cell(row_idx, 2).text = f"{row.ret_20d:.2f}"
        table.cell(row_idx, 3).text = f"{row.momentum63:.2f}"
        table.cell(row_idx, 4).text = f"{row.momentum126:.2f}"
        table.cell(row_idx, 5).text = f"{row.rsi14:.2f}"
        table.cell(row_idx, 6).text = str(row.ma_stack)
    _format_table_font(table, size=10, header=True)
    for row_idx in range(1, 12):
        for col_idx in [1, 2, 3, 4]:
            value = float(table.cell(row_idx, col_idx).text)
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            if value >= 0:
                cell.fill.fore_color.rgb = RGBColor(226, 244, 230)
            else:
                cell.fill.fore_color.rgb = RGBColor(251, 231, 231)
    _add_note(slide, "這頁是前 11 名趨勢體質檢查，找出短線回檔但中長線仍強的標的。")

    # 9. Top11 multi-factor profile
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(slide, RGBColor(255, 255, 255))
    _add_title(slide, "前 11 名：因子拆解", "Trend / Momentum / Value / Fundamental / Risk")
    top11_factor = top11[["symbol", "name", "trend_score", "momentum_score", "value_score", "fundamental_score", "risk_control_score"]]
    t = slide.shapes.add_table(12, 7, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.6)).table
    headers = ["股票", "Trend", "Momentum", "Value", "Fundamental", "Risk", "總評"]
    for i, h in enumerate(headers):
        t.cell(0, i).text = h
    for idx, row in enumerate(top11_factor.itertuples(index=False), start=1):
        t.cell(idx, 0).text = f"{int(row.symbol)} {row.name}"
        vals = [row.trend_score, row.momentum_score, row.value_score, row.fundamental_score, row.risk_control_score]
        for col, val in enumerate(vals, start=1):
            t.cell(idx, col).text = f"{float(val):.1f}"
        momentum = float(row.momentum_score)
        value = float(row.value_score)
        if momentum >= 80 and value >= 50:
            rating = "強趨勢+中性估值"
        elif momentum >= 80 and value < 50:
            rating = "強趨勢+偏貴"
        elif momentum < 80 and value >= 60:
            rating = "修正後價值"
        else:
            rating = "混合訊號"
        t.cell(idx, 6).text = rating
    _format_table_font(t, size=10, header=True)
    _add_note(slide, "不要只看總分。這頁用來分辨高分是靠哪個因子推上去。")

    # 10. Risk and scenario
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(slide, PALETTE["bg_light"])
    _add_title(slide, "風險面與情境推演", "波動、RSI、均線結構的警示條件")
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.8), Inches(12.0), Inches(4.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(232, 240, 250)
    box.line.color.rgb = RGBColor(200, 214, 233)
    tf = box.text_frame
    tf.clear()
    scenarios = [
        f"1) 基準情境：63D/126D 平均維持正值（{summary['avg_returns']['63D']:.2f}% / {summary['avg_returns']['126D']:.2f}%），採分批偏多。",
        f"2) 過熱情境：若 Top11 RSI 中位數 > 75（目前 {top11['rsi14'].median():.2f}），降低追價倉位。",
        f"3) 風險放大：若 Top11 波動中位數 > {top11['volatility20'].median():.2f}%，單檔風險 budget 下修。",
        "4) 失效條件：MA 結構由 上/上/上 轉為 下/下/下 且 20D 報酬轉負，視為結構轉弱。",
        "5) 執行原則：優先持有總分高且 Momentum 與 Fundamental 同步高分的標的。",
    ]
    for idx, line in enumerate(scenarios):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Calibri"
        p.font.size = Pt(17 if idx == 0 else 15)
        p.font.bold = idx == 0
        p.font.color.rgb = PALETTE["ink_dark"]
    _add_note(slide, "這頁把分數翻成交易語言：什麼時候加碼、什麼時候降風險。")

    # 11. Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(slide, PALETTE["bg_dark"])
    _add_title(slide, "結論與下一步", "把訊號變成可執行監控清單", dark=True)
    close_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.8), Inches(3.9))
    ctf = close_box.text_frame
    ctf.clear()
    points = [
        "AI 主題整體仍是中長期偏多結構，短線波動偏高。",
        "前 11 名中，優先看 Momentum + Fundamental 同步高分且 MA 結構完整者。",
        "每週更新 5D/20D 勝率與 Top11 因子拆解，避免只看總分。",
        "若要進一步提升穩健度，下一步可加入歷史滾動回測與權重敏感度分析。",
    ]
    for idx, text in enumerate(points):
        p = ctf.paragraphs[0] if idx == 0 else ctf.add_paragraph()
        p.text = f"- {text}"
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(232, 242, 252)
    _add_note(slide, "收尾強調：這份報告是決策支架，不是自動下單系統。")

    prs.save(path)


def main() -> int:
    args = parse_args()
    csv_path = Path(args.csv_path)
    if not csv_path.exists():
        print(f"[ai-ppt] error: csv not found: {csv_path}")
        return 1
    as_of = datetime.strptime(args.as_of, "%Y-%m-%d").date()
    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    df = pd.read_csv(csv_path)
    numeric_cols = [
        "rank",
        "total_score",
        "trend_score",
        "momentum_score",
        "value_score",
        "fundamental_score",
        "risk_control_score",
        "ret_5d",
        "ret_20d",
        "momentum63",
        "momentum126",
        "rsi14",
        "volatility20",
    ]
    for col in numeric_cols:
        df[col] = _to_num(df[col])
    df = df.sort_values(["rank", "total_score"], ascending=[True, False]).reset_index(drop=True)
    top11 = df.head(11).copy()
    summary = _build_summary(df)

    md_path = output_dir / f"ai-sector-visual-analysis-{as_of.strftime('%Y%m%d')}.md"
    ppt_path = output_dir / f"ai-sector-briefing-{as_of.strftime('%Y%m%d')}.pptx"
    _write_markdown(md_path, as_of, df, top11, summary)
    _build_ppt(ppt_path, as_of, df, top11, summary)

    # lightweight csv export for top11 table reuse
    top11_csv = output_dir / f"ai-top11-{as_of.strftime('%Y%m%d')}.csv"
    with top11_csv.open("w", encoding="utf-8-sig", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(["rank", "symbol", "name", "total_score", "ret_5d", "ret_20d", "momentum63", "momentum126", "rsi14", "volatility20", "ma_stack"])
        for _, row in top11.iterrows():
            writer.writerow(
                [
                    int(row["rank"]),
                    row["symbol"],
                    row["name"],
                    f"{float(row['total_score']):.2f}",
                    f"{float(row['ret_5d']):.2f}",
                    f"{float(row['ret_20d']):.2f}",
                    f"{float(row['momentum63']):.2f}",
                    f"{float(row['momentum126']):.2f}",
                    f"{float(row['rsi14']):.2f}",
                    f"{float(row['volatility20']):.2f}",
                    row["ma_stack"],
                ]
            )

    print(f"[ai-ppt] analysis: {md_path}")
    print(f"[ai-ppt] deck: {ppt_path}")
    print(f"[ai-ppt] top11: {top11_csv}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
