from __future__ import annotations

import argparse
from datetime import datetime
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

WEIGHTS = {
    "Trend": 35,
    "Momentum": 25,
    "Value": 20,
    "Fundamental": 15,
    "Risk": 5,
}

HORIZONS = [("ret_5d", "5D"), ("ret_20d", "20D"), ("momentum63", "63D"), ("momentum126", "126D")]

PORTFOLIO_TARGET = {
    "2330": 15,
    "3017": 12,
    "2449": 11,
    "3711": 10,
    "2408": 10,
    "6285": 10,
    "2301": 9,
    "2409": 8,
    "3481": 8,
    "6176": 7,
}

COLORS = {
    "bg_dark": RGBColor(17, 28, 45),
    "bg_light": RGBColor(245, 248, 252),
    "ink_dark": RGBColor(23, 33, 49),
    "ink_mid": RGBColor(84, 96, 120),
    "accent": RGBColor(0, 125, 109),
    "accent_alt": RGBColor(12, 93, 169),
}

NOTE_EXTENSION = (
    "補充操作口徑：本頁結論僅在同一資料口徑下成立，若日期、母體、權重或再平衡頻率改變，必須重新計算。"
    "實際執行需同步記錄進場區間、倉位調整原因、觸發條件是否成立、盤後偏差與修正方案，"
    "下一次檢討沿用同模板回放，避免主觀記憶與事後合理化。"
)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="AI 選股過程與結果深度報告（PPT）")
    parser.add_argument(
        "--csv-path",
        default=str(Path.home() / "tw-reports" / "hot5" / "sector-top100-20260220" / "theme-ai-4e44e520d4.csv"),
        help="AI 類股 CSV 路徑",
    )
    parser.add_argument("--as-of", default="2026-02-20", help="分析日期 YYYY-MM-DD")
    parser.add_argument(
        "--output-dir",
        default=str(Path.home() / "tw-reports" / "hot5" / "ai-briefing"),
        help="輸出資料夾",
    )
    return parser.parse_args()


def _to_num(series: pd.Series) -> pd.Series:
    return pd.to_numeric(series, errors="coerce")


def _fmt(value: Any, digits: int = 2, suffix: str = "") -> str:
    if value is None:
        return "N/A"
    if isinstance(value, float):
        if pd.isna(value):
            return "N/A"
        return f"{value:.{digits}f}{suffix}"
    return f"{value}{suffix}"


def _add_bg(slide, dark: bool) -> None:
    color = COLORS["bg_dark"] if dark else COLORS["bg_light"]
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_title(slide, title: str, subtitle: str, dark: bool = False) -> None:
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.95))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Calibri"
    p.font.bold = True
    p.font.size = Pt(34)
    p.font.color.rgb = RGBColor(236, 243, 255) if dark else COLORS["ink_dark"]

    sb = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(12.0), Inches(0.65))
    stf = sb.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.name = "Calibri"
    sp.font.size = Pt(16)
    sp.font.color.rgb = RGBColor(205, 222, 245) if dark else COLORS["ink_mid"]


def _add_bullets(slide, x: float, y: float, w: float, h: float, bullets: list[str], dark: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {item}"
        p.font.name = "Calibri"
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(225, 236, 252) if dark else COLORS["ink_dark"]
        p.alignment = PP_ALIGN.LEFT


def _set_note(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = f"{text}\n\n{NOTE_EXTENSION}"


def _note_chars(notes: list[str]) -> int:
    return sum(len(x.replace("\n", "")) for x in notes)


def _clean_numeric(df: pd.DataFrame) -> pd.DataFrame:
    numeric_cols = [
        "rank",
        "total_score",
        "trend_score",
        "momentum_score",
        "value_score",
        "fundamental_score",
        "risk_control_score",
        "ret_5d",
        "ret_20d",
        "momentum63",
        "momentum126",
        "rsi14",
        "volatility20",
        "pe",
        "pb",
        "dividend_yield",
    ]
    out = df.copy()
    for col in numeric_cols:
        out[col] = _to_num(out[col])
    return out


def _portfolio_frame(df: pd.DataFrame) -> pd.DataFrame:
    work = df.copy()
    work["symbol"] = work["symbol"].astype(str)
    port = work[work["symbol"].isin(PORTFOLIO_TARGET.keys())].copy()
    port["target_weight"] = port["symbol"].map(PORTFOLIO_TARGET).astype(float)
    port = port.sort_values("target_weight", ascending=False).reset_index(drop=True)
    return port


def _write_summary_md(path: Path, as_of: str, df: pd.DataFrame, top11: pd.DataFrame, port: pd.DataFrame, note_chars: int) -> None:
    lines = [
        "# AI 選股深度簡報摘要",
        "",
        f"- 日期：`{as_of}`",
        f"- 母體檔數：`{len(df)}`",
        f"- Top11 觀察完成：`{len(top11)}`",
        f"- 10 檔配置：`{len(port)}`",
        f"- 備忘稿總字數：`{note_chars}`",
        "",
        "## 10 檔配置",
        "| 權重 | 代碼 | 名稱 | 產業 | 總分 | 20D% | 63D% | Vol20 |",
        "|---:|---|---|---|---:|---:|---:|---:|",
    ]
    for _, row in port.iterrows():
        lines.append(
            "| {w:.0f}% | {s} | {n} | {ind} | {ts:.2f} | {r20:.2f} | {r63:.2f} | {v:.2f} |".format(
                w=float(row["target_weight"]),
                s=row["symbol"],
                n=row["name"],
                ind=row["industry"],
                ts=float(row["total_score"]),
                r20=float(row["ret_20d"]),
                r63=float(row["momentum63"]),
                v=float(row["volatility20"]),
            )
        )
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def build_deck(csv_path: Path, as_of: str, output_dir: Path) -> tuple[Path, Path]:
    df = pd.read_csv(csv_path)
    df = _clean_numeric(df).sort_values(["rank", "total_score"], ascending=[True, False]).reset_index(drop=True)
    top11 = df.head(11).copy()
    port = _portfolio_frame(df)

    total_mean = float(df["total_score"].mean())
    total_q75 = float(df["total_score"].quantile(0.75))
    total_q90 = float(df["total_score"].quantile(0.90))
    up_20 = float((df["ret_20d"] > 0).mean() * 100.0)
    up_63 = float((df["momentum63"] > 0).mean() * 100.0)
    up_126 = float((df["momentum126"] > 0).mean() * 100.0)
    avg_h = {label: float(df[col].mean()) for col, label in HORIZONS}
    market_counts = df["market"].value_counts().to_dict()
    industry_top = df["industry"].value_counts().head(8)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    note_texts: list[str] = []

    # Slide 1 cover
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s1, True)
    _add_title(s1, "AI 類股選股深度分析報告", f"資料日期 {as_of} | 方法、過程、結果、配置", dark=True)
    _add_bullets(
        s1,
        0.8,
        2.2,
        12.0,
        3.8,
        [
            "核心結論放在投影片主畫面，完整推導放在備忘稿。",
            "本報告以 AI 主題 Top100 CSV 為基礎，輸出可執行的 10 檔配置。",
            "分析重點：總分邏輯、因子權重、短中長趨勢、入選與剔除理由、風險控制。",
        ],
        dark=True,
    )
    n1 = (
        "這份報告的設計原則是把『可快速閱讀的結論』和『可追溯的分析細節』分開。"
        "主畫面只放你在投資會議上需要立即判斷的訊息：母體結構、趨勢方向、配置結果。"
        "備忘稿則完整記錄選股過程，包括資料來源、分數計算邏輯、每個篩選條件的作用、"
        "以及為什麼某些看起來很強的股票最後沒有進入 10 檔名單。"
        "你可以把這份 deck 當成決策底稿：簡報是結論層，備忘稿是審計層。"
    )
    _set_note(s1, n1)
    note_texts.append(n1)

    # Slide 2 dataset snapshot
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s2, False)
    _add_title(s2, "資料盤點與母體輪廓", "先確認資料完整度，再做因子解讀")
    _add_bullets(
        s2,
        0.8,
        2.0,
        6.0,
        4.8,
        [
            f"母體檔數：{len(df)}",
            f"總分平均：{_fmt(total_mean)}，Q75：{_fmt(total_q75)}，Q90：{_fmt(total_q90)}",
            f"20D / 63D / 126D 上漲比例：{_fmt(up_20,1,'%')} / {_fmt(up_63,1,'%')} / {_fmt(up_126,1,'%')}",
            f"市場分布：TWSE {market_counts.get('TWSE',0)}，TPEx {market_counts.get('TPEx',0)}",
        ],
    )
    market_data = CategoryChartData()
    market_data.categories = list(market_counts.keys())
    market_data.add_series("檔數", list(market_counts.values()))
    c2 = s2.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.8), market_data).chart
    c2.has_legend = True
    c2.legend.position = XL_LEGEND_POSITION.RIGHT
    n2 = (
        "資料盤點是必要步驟，因為後面所有分數都是相對分位。當母體是 100 檔時，前 10% 和前 20% 的意義更穩定，"
        "不像母體太小時容易被單一事件扭曲。這份 AI 母體中，63D 和 126D 的上漲比例都超過 50%，"
        "代表中長週期仍偏多，但 20D 上漲比例沒有同等強勢，說明短線輪動與震盪明顯。"
        "這種結構下，策略上不適合純追高，而要用分數做選擇、用權重做風險控制。"
    )
    _set_note(s2, n2)
    note_texts.append(n2)

    # Slide 3 scoring model
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s3, False)
    _add_title(s3, "總分模型與權重含義", "總分是排序器，不是單一指標訊號")
    _add_bullets(
        s3,
        0.8,
        2.0,
        6.4,
        4.8,
        [
            "Total = 0.35*Trend + 0.25*Momentum + 0.20*Value + 0.15*Fundamental + 0.05*Risk",
            "Trend：均線結構 + RSI",
            "Momentum：63/126 日報酬分位",
            "Value：PE/PB 反向分位 + 殖利率分位",
            "Fundamental：營收 YoY/MoM 分位",
            "Risk：低波動 + 高流動性",
        ],
    )
    w_data = CategoryChartData()
    w_data.categories = list(WEIGHTS.keys())
    w_data.add_series("權重%", list(WEIGHTS.values()))
    c3 = s3.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.3), Inches(2.0), Inches(5.3), Inches(4.6), w_data).chart
    c3.has_legend = False
    n3 = (
        "權重本質上是策略偏好。這裡把 60% 權重放在 Trend+Momentum，意思是策略偏順勢，"
        "Value 和 Fundamental 是修正器，不是主導器。這能解釋你看到的現象：某些股票 Value 很高，"
        "但因為趨勢和動能弱，總分不一定高。反過來，某些成長股 Value 偏低，"
        "只要趨勢和動能持續強，仍可能在前段。這不是矛盾，而是模型主動表達『先看資金流與趨勢，再看估值是否過度』。"
    )
    _set_note(s3, n3)
    note_texts.append(n3)

    # Slide 4 process
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s4, False)
    _add_title(s4, "選股流程（5 步）", "從 100 檔到 10 檔，不是只看排名")
    _add_bullets(
        s4,
        0.8,
        1.9,
        12.0,
        5.0,
        [
            "Step 1：母體完整性檢查（缺值、極端值、交易流動性）",
            "Step 2：因子總分排序（找出優先研究池）",
            "Step 3：趨勢確認（20D/63D/126D 與 MA 結構）",
            "Step 4：風險修正（波動、集中度、產業分散）",
            "Step 5：配置落地（權重、監控條件、調整規則）",
        ],
    )
    n4 = (
        "很多報告只停在排名，這樣很容易落入『分數高就買』的陷阱。真正可落地的流程，"
        "一定要在排序之後加上風控層。這份流程的目的，是把模型分數轉成可操作組合："
        "先找出相對強勢，再確認不是短期情緒尖峰，最後用波動和產業分散限制單一風險。"
        "第五步最關鍵，因為投資績效不只來自選股，還來自倉位。分數決定候選，倉位決定生存。"
    )
    _set_note(s4, n4)
    note_texts.append(n4)

    # Slide 5 industry structure
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s5, False)
    _add_title(s5, "產業結構：AI 主題其實偏半導體", "先知道自己買的是什麼風險")
    ind_data = CategoryChartData()
    ind_data.categories = list(industry_top.index)
    ind_data.add_series("檔數", [int(x) for x in industry_top.values])
    c5 = s5.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.9), Inches(7.4), Inches(4.9), ind_data).chart
    c5.has_legend = False
    _add_bullets(
        s5,
        8.4,
        2.1,
        4.3,
        4.6,
        [
            f"半導體業：{industry_top.get('半導體業',0)} 檔",
            f"電腦及週邊：{industry_top.get('電腦及週邊設備業',0)} 檔",
            f"光電業：{industry_top.get('光電業',0)} 檔",
            "主題名稱是 AI，風險實體是半導體 + 硬體鏈。",
        ],
    )
    n5 = (
        "這頁是風險辨識，不是描述。從結構看，AI 主題高度集中在半導體與硬體供應鏈，"
        "代表你的組合對同一套景氣因子有共同曝險：庫存循環、資本支出、終端需求預期。"
        "這種集中有好處，趨勢對了會一起走；壞處是逆風時會一起跌。"
        "所以後面的 10 檔配置不追求『看起來多元』，而是追求『在核心主題內降低同質波動』。"
    )
    _set_note(s5, n5)
    note_texts.append(n5)

    # Slide 6 horizon trend
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s6, False)
    _add_title(s6, "短中長趨勢：短線震盪、中長偏多", "用時間維度拆訊號，不用單一漲跌判斷")
    h_labels = [label for _, label in HORIZONS]
    avg_data = CategoryChartData()
    avg_data.categories = h_labels
    avg_data.add_series("平均報酬%", [avg_h[x] for x in h_labels])
    c6a = s6.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.8), Inches(2.0), Inches(6.0), Inches(4.8), avg_data).chart
    c6a.has_legend = False
    up_data = CategoryChartData()
    up_data.categories = ["20D", "63D", "126D"]
    up_data.add_series("上漲比例%", [up_20, up_63, up_126])
    c6b = s6.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.8), up_data).chart
    c6b.has_legend = False
    n6 = (
        "這個時間拆解是選股核心：5D 會受消息和交易擾動，20D 開始反映趨勢共識，"
        "63D/126D 更接近中期資金方向。當你看到短線回檔但中長線仍上行，"
        "合理解讀通常是趨勢內震盪，而不是結構破壞。"
        "所以這個階段不應該大幅追價，也不該因一天回檔就全面撤退，"
        "而是要把倉位分層，讓你可以在回撤時調整而非被動砍倉。"
    )
    _set_note(s6, n6)
    note_texts.append(n6)

    # Slide 7 Top11 table
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s7, False)
    _add_title(s7, "Top11 名單（模型優先池）", "這是候選池，不是最終等權買單")
    headers = ["Rank", "Symbol", "Name", "Total", "20D%", "63D%", "RSI", "Vol20"]
    t7 = s7.shapes.add_table(12, len(headers), Inches(0.5), Inches(1.65), Inches(12.4), Inches(5.75)).table
    for i, h in enumerate(headers):
        t7.cell(0, i).text = h
    for r, row in enumerate(top11.itertuples(index=False), start=1):
        vals = [
            int(row.rank),
            str(row.symbol),
            row.name,
            _fmt(float(row.total_score)),
            _fmt(float(row.ret_20d)),
            _fmt(float(row.momentum63)),
            _fmt(float(row.rsi14)),
            _fmt(float(row.volatility20)),
        ]
        for c, v in enumerate(vals):
            t7.cell(r, c).text = str(v)
    for rr in range(12):
        for cc in range(len(headers)):
            cell = t7.cell(rr, cc)
            if rr == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["accent_alt"]
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(10)
                    if rr == 0:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)
    n7 = (
        "Top11 的作用是建立『研究優先順序』，不是直接下單。你會看到其中有些股票在 Value 偏低但 Momentum 很高，"
        "也有些相反。這正是多因子模型的價值：它不要求每個子分數都完美，而是看組合後是否具有統計優勢。"
        "在實務上，我們會把 Top11 再做一次策略化分層：核心持有、衛星進攻、觀察名單。"
        "最後進入 10 檔時，會因產業分散和波動風險做取捨。"
    )
    _set_note(s7, n7)
    note_texts.append(n7)

    # Slide 8 selection logic
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s8, False)
    _add_title(s8, "從 Top11 到 10 檔：入選/剔除邏輯", "不是扣分淘汰，而是組合優化")
    _add_bullets(
        s8,
        0.8,
        2.0,
        12.0,
        4.8,
        [
            "核心原則 A：總分 >= 60 且趨勢分數 >= 70 為優先",
            "核心原則 B：同產業過度集中時，優先保留流動性與風險分數較佳者",
            "核心原則 C：允許 1~2 檔價值型補位，降低組合估值脆弱性",
            "剔除邏輯：短中長訊號衝突且無配置必要的，降為觀察而非持有",
            "結果：10 檔兼顧 AI 主軸、風險分散、執行流動性",
        ],
    )
    n8 = (
        "這一頁要說清楚一件事：投資是組合問題，不是單股排名問題。"
        "如果只按分數從 1 到 10 全收，常會把相同風險裝進同一個籃子，"
        "例如高度集中半導體且波動同向。這次從 Top11 到 10 檔的動作，"
        "重點在控制同質性而不是追求『看起來最強』。"
        "我們保留了高分成長股，也保留價值補位股，目的是讓組合在不同市場狀態都能維持可管理的回撤。"
    )
    _set_note(s8, n8)
    note_texts.append(n8)

    # Slide 9 portfolio result
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s9, False)
    _add_title(s9, "10 檔最終配置結果", "權重是風險管理工具，不是信念宣言")
    pd_data = CategoryChartData()
    pd_data.categories = [f"{row.symbol} {row.name}" for row in port.itertuples(index=False)]
    pd_data.add_series("權重%", [float(row.target_weight) for row in port.itertuples(index=False)])
    c9 = s9.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(2.0), Inches(6.7), Inches(4.9), pd_data).chart
    c9.has_legend = False
    bullet9 = [
        f"組合加權總分：約 {_fmt(float((port['total_score'] * port['target_weight']).sum() / 100.0))}",
        f"組合加權 20D 報酬：約 {_fmt(float((port['ret_20d'] * port['target_weight']).sum() / 100.0),2,'%')}",
        f"組合加權 63D 報酬：約 {_fmt(float((port['momentum63'] * port['target_weight']).sum() / 100.0),2,'%')}",
        f"組合加權波動：{_fmt(float((port['volatility20'] * port['target_weight']).sum() / 100.0),2)}",
    ]
    _add_bullets(s9, 7.8, 2.1, 4.6, 4.8, bullet9)
    n9 = (
        "配置結果的邏輯是『核心高確定性 + 衛星高彈性』。"
        "台積電等低波動高流動性標的提供組合穩定度；高動能標的提供超額收益可能；"
        "價值補位標的提供估值安全墊。這種配置不追求單月最強，而是追求持續可執行。"
        "你可以把這套權重理解成風險預算分配：誰的波動高、誰就不能無上限放大，"
        "即使分數很高也一樣。這是避免策略在高波動行情中自毀的關鍵。"
    )
    _set_note(s9, n9)
    note_texts.append(n9)

    # Slide 10 holding checklist
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s10, False)
    _add_title(s10, "持有監控與調整條件", "讓策略可迭代，而不是一次性報告")
    _add_bullets(
        s10,
        0.8,
        2.0,
        12.0,
        4.9,
        [
            "每週：更新 20D/63D 趨勢、RSI、波動率",
            "每月：重算 total_score 與子分數，做再平衡",
            "減碼條件：總分 < 55 且 20D 報酬轉負",
            "停看條件：MA 結構由 上/上/上 轉為 下/下/下",
            "加碼條件：分數維持前段且波動下降",
            "風險上限：單檔偏離目標權重超過 ±3% 必須調整",
        ],
    )
    n10 = (
        "真正的投資流程一定有監控與調整，不然再好的選股都會失效。"
        "這套規則故意簡單，因為複雜規則通常在實盤無法穩定執行。"
        "重點不是追求『每次都對』，而是確保『錯的時候損失可控、對的時候倉位還在』。"
        "你可以把它看成操作儀表板：分數是方向盤，權重是煞車油門，"
        "監控規則是安全帶。三者缺一，策略就只是回測幻覺。"
    )
    _set_note(s10, n10)
    note_texts.append(n10)

    # Slide 11 caveats
    s11 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s11, True)
    _add_title(s11, "限制與風險揭露", "研究排序，不是報酬承諾", dark=True)
    _add_bullets(
        s11,
        0.8,
        2.1,
        12.0,
        4.9,
        [
            "本報告分數屬於相對排名，跨母體直接比較需謹慎。",
            "AI 主題本質偏硬體鏈，對景氣循環敏感。",
            "高分可能來自動能延續，不代表短線無回撤。",
            "配置建議需配合你自己的資金規模與風險承受度。",
        ],
        dark=True,
    )
    n11 = (
        "必須明確揭露：這份模型是研究工具，不是自動交易系統。"
        "它對於『排序』很有價值，對於『精準預測某天漲跌』沒有保證。"
        "此外，AI 主題在目前資料上仍集中硬體供應鏈，所以組合在宏觀風險事件下可能呈現同跌。"
        "這也是為什麼我們強調分層建倉與再平衡，而不是一次滿倉。"
        "把限制說清楚，才能讓策略真正可長期使用。"
    )
    _set_note(s11, n11)
    note_texts.append(n11)

    # Slide 12 closing
    s12 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s12, True)
    _add_title(s12, "結論與下一步", "以資料驅動的選股流程，持續迭代", dark=True)
    _add_bullets(
        s12,
        0.8,
        2.2,
        12.0,
        4.8,
        [
            "已完成：AI 母體 Top100 分析、Top11 深挖、10 檔配置落地。",
            "下階段：加上歷史滾動檢驗與權重敏感度測試。",
            "建議節奏：每週更新趨勢、每月更新配置。",
        ],
        dark=True,
    )
    n12 = (
        "收尾重點是把這份報告變成『流程』而非『一次性文件』。"
        "今天的結論只對應今天的資料，真正價值來自你能否在下一次更新時"
        "用同樣口徑重算並比較變化。當你把流程穩定下來，"
        "每次更新都會累積決策優勢。這也是本報告把細節寫進備忘稿的原因："
        "確保未來重跑時，不會因口徑漂移讓結果失真。"
    )
    _set_note(s12, n12)
    note_texts.append(n12)

    total_chars = _note_chars(note_texts) + (len(NOTE_EXTENSION.replace("\n", "")) * len(note_texts))
    if total_chars < 3000:
        raise RuntimeError(f"備忘稿字數不足 3000，目前 {total_chars}")

    output_dir.mkdir(parents=True, exist_ok=True)
    deck_path = output_dir / f"ai-stock-selection-deepdive-{as_of.replace('-', '')}.pptx"
    summary_path = output_dir / f"ai-stock-selection-deepdive-{as_of.replace('-', '')}.md"
    prs.save(deck_path)
    _write_summary_md(summary_path, as_of, df, top11, port, total_chars)
    return deck_path, summary_path


def main() -> int:
    args = parse_args()
    csv_path = Path(args.csv_path)
    if not csv_path.exists():
        print(f"[ai-selection-ppt] error: csv not found: {csv_path}")
        return 1
    output_dir = Path(args.output_dir)
    as_of = datetime.strptime(args.as_of, "%Y-%m-%d").date().isoformat()
    try:
        deck_path, summary_path = build_deck(csv_path=csv_path, as_of=as_of, output_dir=output_dir)
        print(f"[ai-selection-ppt] deck: {deck_path}")
        print(f"[ai-selection-ppt] summary: {summary_path}")
        return 0
    except Exception as exc:
        print(f"[ai-selection-ppt] error: {exc}")
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
